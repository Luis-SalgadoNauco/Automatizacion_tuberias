from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime


def agregar_slide(prs, titulo, contenido):
    """Agregar slide estándar con título y contenido"""
    slide_layout = prs.slide_layouts[1]  # Título + contenido
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = titulo
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    if isinstance(contenido, dict):
        for k, v in contenido.items():
            p = text_frame.add_paragraph()
            p.text = f"{k}: {v}"
            p.level = 1
    elif isinstance(contenido, list):
        for item in contenido:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 1
    else:
        text_frame.text = contenido


def generar_presentacion_ejecutiva():
    prs = Presentation()

    executive_summary = {
        'titulo': 'Implementación de Pipeline ETL para Analytics E-commerce',

        'resumen_ejecutivo': (
            "Se implementó un pipeline ETL moderno que transforma datos crudos "
            "de e-commerce en insights accionables, reduciendo tiempos operativos "
            "y mejorando la toma de decisiones."
        ),

        'problema': [
            "Reportes manuales tomaban hasta 3 días",
            "Datos inconsistentes entre sistemas",
            "Falta de visibilidad en tiempo casi real",
            "Altos costos operativos"
        ],

        'solucion': [
            "Pipeline ETL automatizado",
            "Extracción desde múltiples fuentes",
            "Data Warehouse dimensional",
            "Dashboards self-service"
        ],

        'beneficios': {
            "Eficiencia": "Reportes de 3 días a 4 horas",
            "Calidad": "99.5% de datos validados",
            "Escalabilidad": "Soporta 10x crecimiento",
            "ROI": "Recuperación en 8 meses"
        },

        'metricas': {
            "Volumen procesado": "500GB/día",
            "Latencia": "< 30 minutos",
            "Disponibilidad": "99.9%",
            "Usuarios": "150+ analistas"
        },

        'riesgos_mitigados': [
            "Monitoreo 24/7",
            "Backups diarios",
            "Arquitectura tolerante a fallos",
            "Testing automatizado"
        ],

        'roadmap': {
            "Fase 1": "Pipeline core operativo",
            "Fase 2": "ML y analytics avanzados",
            "Fase 3": "Personalización en tiempo real"
        },

        'recomendaciones': [
            "Expandir uso a más áreas",
            "Incorporar modelos predictivos",
            "Integrar CRM",
            "Capacitar nuevos usuarios"
        ]
    }

    # Slide título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = executive_summary['titulo']
    slide.placeholders[1].text = f"Presentación ejecutiva\n{datetime.now().strftime('%Y-%m-%d')}"

    # Slides de contenido
    agregar_slide(prs, "Resumen Ejecutivo", executive_summary['resumen_ejecutivo'])
    agregar_slide(prs, "El Problema", executive_summary['problema'])
    agregar_slide(prs, "La Solución", executive_summary['solucion'])
    agregar_slide(prs, "Beneficios Clave", executive_summary['beneficios'])
    agregar_slide(prs, "Métricas de Éxito", executive_summary['metricas'])
    agregar_slide(prs, "Riesgos Mitigados", executive_summary['riesgos_mitigados'])
    agregar_slide(prs, "Roadmap", executive_summary['roadmap'])
    agregar_slide(prs, "Recomendaciones", executive_summary['recomendaciones'])

    output_file = "presentacion_ejecutiva_dia4.pptx"
    prs.save(output_file)

    print(f"Presentación ejecutiva generada: {output_file}")


if __name__ == "__main__":
    generar_presentacion_ejecutiva()
